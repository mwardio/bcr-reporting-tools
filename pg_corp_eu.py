import pandas as pd
from io import BytesIO
import requests as r
from datetime import datetime
from pptx import Presentation, parts
from pptx.chart.data import CategoryChartData
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.parts.image import Image
from pptx.oxml.ns import qn
from helpers.pg_helper import *




def pg_corp_eu_report(region, fiscal_period, fiscal_dates_df, regions_df, queries_df):

    bcr_token = "d449d539-f164-4f88-b6a7-3df75b54c44a"
    headers = {"authorization": f"bearer {bcr_token}", "content-type": "application/json"}

    project_id = 1998383733
    query_group_id = 2048563524
    exclusion_categories = 26995986


    region_id = regions_df.loc[regions_df['country'] == region, 'id'].item()
    report_countries = regions_df.loc[regions_df['country'] == region, 'reportRegion'].item()

    start_date = fiscal_dates_df.loc[fiscal_dates_df['fiscal_period'] == fiscal_period, 'start_date'].item()
    end_date = fiscal_dates_df.loc[fiscal_dates_df['fiscal_period'] == fiscal_period, 'end_date'].item()
    prev_start = fiscal_dates_df.loc[fiscal_dates_df['fiscal_period'] == fiscal_period, 'prev_start'].item()
    prev_end = fiscal_dates_df.loc[fiscal_dates_df['fiscal_period'] == fiscal_period, 'prev_end'].item()


    ## PULL P&G DAILY VOLUME BY SENTIMENT

    pg_daily_sentiment_call = f'https://api.brandwatch.com/projects/{project_id}/data/volume/days/sentiment'
    params = {'queryId': 2002772325,
            'startDate': start_date,
            'endDate': f"{end_date}T23:59:59",
            'location': region_id,
            'xcategory': exclusion_categories,}
            #'timezone': "America/New_York"}
    pg_daily_sentiment_chart = r.get(pg_daily_sentiment_call, params=params, headers=headers)

    pg_daily_sentiment = []
    for x in pg_daily_sentiment_chart.json()['results']:
        #date = datetime.fromisoformat(x['id']).astimezone(pytz.timezone('America/New_York')).strftime("%m/%#d/%Y")
        date = datetime.fromisoformat(x['id']).strftime("%#d-%b")
        data = {"date": date}
        for y in x['values']:
            value_data = {y['id']: y['value']}
            data.update(value_data)
        pg_daily_sentiment.append(data)

    pg_daily_sent_df = pd.DataFrame(pg_daily_sentiment)


    ## PULL P&G VOLUME BY MEDIA TYPE

    pg_media_volume_call = f'https://api.brandwatch.com/projects/{project_id}/data/volume/queries/pageTypes'
    params = {'queryId': 2002772325,
            'startDate': start_date,
            'endDate': f"{end_date}T23:59:59",
            'location': region_id,
            'xcategory': exclusion_categories,}
            #'timezone': "America/New_York"}
    pg_media_volume_chart = r.get(pg_media_volume_call, params=params, headers=headers)

    pg_media_volume_data = []
    for x in pg_media_volume_chart.json()['results']:
        for y in x['values']:
            pg_media_volume_data.append({"id": y['id'], "name": y['name'], "mentions": y['value'],})
        
    pg_media_volume_df = pd.DataFrame(pg_media_volume_data).sort_values(by='mentions', ascending=False, ignore_index=True)
    pg_media_volume_df['share'] = round((pg_media_volume_df['mentions'] / pg_media_volume_df['mentions'].sum()) * 100, 1)


    ## DETERMINE P&G TOTAL BRAND VOLUME

    pg_total_mentions = pg_media_volume_df['mentions'].sum()


    ## PULL VOLUME BY SENTIMENT BY BRAND/QUERY (ALL)

    brands_volume_call = f'https://api.brandwatch.com/projects/{project_id}/data/volume/queries/sentiment'
    params = {#'queryGroupId': query_group_id,
            'queryId': [2002772325,2002797072,2002797068,2002797067,2002797070,2002797217,2002797034,2002797071],
            'startDate': start_date,
            'endDate': f"{end_date}T23:59:59",
            'location': region_id,
            'xcategory': exclusion_categories,}
            #'timezone': "America/New_York"}
    brands_volume_chart = r.get(brands_volume_call, params=params, headers=headers)

    brands_volume_data = []
    for x in brands_volume_chart.json()['results']:
        data = {"queryId": x['id'], "queryName": x['name']}
        for y in x['values']:
            value_data = {y['id']: y['value']}
            data.update(value_data)
        brands_volume_data.append(data)

    brands_volume_df = pd.DataFrame(brands_volume_data)
    brands_volume_df['total'] = brands_volume_df['positive']+brands_volume_df['neutral']+brands_volume_df['negative']
    brands_volume_df = brands_volume_df.merge(queries_df[['queryId','shortName']], on='queryId', how='left')
    brands_volume_df.sort_values(by="total", ascending=False, ignore_index=True, inplace=True)

    brands_volume_df['pos%'] = brands_volume_df['positive'] / brands_volume_df['total']
    brands_volume_df['neu%'] = brands_volume_df['neutral'] / brands_volume_df['total']
    brands_volume_df['neg%'] = brands_volume_df['negative'] / brands_volume_df['total']


    ## PULL PREVIOUS PERIOD VOLUME BY BRAND

    prev_period_volume_call = f'https://api.brandwatch.com/projects/{project_id}/data/volume/queryGroups/queries'
    params = {'queryGroupId': query_group_id,
            'queryId': [2002772325,2002797072,2002797068,2002797067,2002797070,2002797217,2002797034,2002797071],
            'startDate': prev_start,
            'endDate': f"{prev_end}T23:59:59",
            'location': region_id,
            'xcategory': exclusion_categories,}
            #'timezone': "America/New_York"}
    prev_period_volume_chart = r.get(prev_period_volume_call, params=params, headers=headers)

    prev_period_volume_data = []
    for x in prev_period_volume_chart.json()['results']:
        for y in x['values']:
            prev_period_volume_data.append({"queryId": y['id'], "queryName": y['name'], "prev_mentions": y['value'],})
        
    prev_period_volume_df = pd.DataFrame(prev_period_volume_data).sort_values(by='prev_mentions', ascending=False, ignore_index=True)
    prev_period_volume_df = prev_period_volume_df.merge(queries_df[['queryId','shortName']], on='queryId', how='left')


    ## ADD 'PREV_MENTIONS' AND '%_CHNG' COLUMNS TO BRANDS_VOLUME_DF

    brands_volume_df = brands_volume_df.merge(prev_period_volume_df[['queryId','prev_mentions']], on='queryId', how='left')
    brands_volume_df['%_chng'] = ((brands_volume_df['total']-brands_volume_df['prev_mentions']) / brands_volume_df['prev_mentions']) * 100


    ## DETERMINE TOP FOUR BRANDS BY VOLUME AND LOOK UP FIGURES

    brand1_name = brands_volume_df.iloc[0,6]
    brand1_volume = int(brands_volume_df.iloc[0,5])
    brand1_prev = int(brands_volume_df.iloc[0,10])
    brand1_chng = int(round(brands_volume_df.iloc[0,11]))

    brand2_name = brands_volume_df.iloc[1,6]
    brand2_volume = int(brands_volume_df.iloc[1,5])
    brand2_prev = int(brands_volume_df.iloc[1,10])
    brand2_chng = int(round(brands_volume_df.iloc[1,11]))

    brand3_name = brands_volume_df.iloc[2,6]
    brand3_volume = int(brands_volume_df.iloc[2,5])
    brand3_prev = int(brands_volume_df.iloc[2,10])
    brand3_chng = int(round(brands_volume_df.iloc[2,11]))

    brand4_name = brands_volume_df.iloc[3,6]
    brand4_volume = int(brands_volume_df.iloc[3,5])
    brand4_prev = int(brands_volume_df.iloc[3,10])
    brand4_chng = int(round(brands_volume_df.iloc[3,11]))


    ## OPEN DECK TEMPLATE AND MAKE INITIAL TEXT REPLACEMENTS

    prs = Presentation("templates/pg_template.pptx")

    replacements = {"start_end_month": f"{datetime.fromisoformat(start_date).strftime("%B %Y")} – {datetime.fromisoformat(end_date).strftime("%B %Y")}",
                    "report_countries": report_countries,
                    "report_region": region,
                    "fiscal_period": fiscal_period}

    for slide in prs.slides:
        for key, value in replacements.items():
            replace_text(slide, key, value)


    ## SLIDE 5(4) - P&G VOLUME + PREV. PERIOD COMPARISON

    pg_prev_volume = int(brands_volume_df.loc[brands_volume_df['queryId'] == 2002772325, 'prev_mentions'].iloc[0])
    pg_vol_chng = round(brands_volume_df.loc[brands_volume_df['queryId'] == 2002772325, '%_chng'].iloc[0])

    replace_text(prs.slides[4],"pg_vol", f"{pg_total_mentions:,}")
    replace_text(prs.slides[4],"pgVol_vs_prev", f"{abs(pg_vol_chng)}% vs. {pg_prev_volume:,} previous period")

    pg_vol_arrow = find_shape_by_name(prs.slides[4], "pg_vol_arrow")
    if pg_vol_chng > 0:
        pg_vol_arrow.fill.fore_color.rgb = RGBColor(151, 215, 0)
        pg_vol_arrow.rotation = 180
    elif pg_vol_chng < 0:
        pg_vol_arrow.fill.fore_color.rgb = RGBColor(255, 109, 86)
        pg_vol_arrow.rotation = 0
    else:
        pg_vol_arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)


    ## SLIDE 5(4) - P&G SENTIMENT PIE CHART

    pg_sent_pivot = {
        'sentiment': ['Positive','Negative','Neutral'],
        'mentions': [int(pg_daily_sent_df['positive'].sum()), int(pg_daily_sent_df['negative'].sum()), int(pg_daily_sent_df['neutral'].sum())]
    }
    pg_sent_df = pd.DataFrame(pg_sent_pivot)


    pg_sent_pie_data = CategoryChartData()
    pg_sent_pie_data.categories = pg_sent_df['sentiment'].to_list()
    pg_sent_pie_data.add_series("mentions", pg_sent_df['mentions'].to_list())

    pg_sentiment_pie = find_shape_by_name(prs.slides[4], "pg_sentiment_pie").chart
    pg_sentiment_pie.replace_data(pg_sent_pie_data)


    ## SLIDE 5(4) - P&G DAILY SENTIMENT CHART

    pg_daily_sent_data = CategoryChartData()
    pg_daily_sent_data.categories = pg_daily_sent_df['date'].to_list()
    pg_daily_sent_data.add_series("positive", pg_daily_sent_df['positive'].to_list())
    pg_daily_sent_data.add_series("neutral", pg_daily_sent_df['neutral'].to_list())
    pg_daily_sent_data.add_series("negative", pg_daily_sent_df['negative'].to_list())

    pg_daily_sent_small = find_shape_by_name(prs.slides[4], "pg_daily_sent_small").chart
    pg_daily_sent_small.replace_data(pg_daily_sent_data)


    ## SLIDE 5(4) - COMPETITIVE SOV CHART

    comp_sov_pie_data = CategoryChartData()
    comp_sov_pie_data.categories = brands_volume_df['shortName'].to_list()
    comp_sov_pie_data.add_series("mentions", brands_volume_df['total'].to_list())

    comp_sov_pie_one = find_shape_by_name(prs.slides[4], "comp_sov_pie_one").chart
    comp_sov_pie_one.replace_data(comp_sov_pie_data)


    ## SLIDE 5(4) - TOP FOUR CHANNEL ICONS, SHARE, VOLUME & LABELS

    # Determine top four channels
    for i in range(4):
        try:
            globals()[f"media{i+1}_name"] = pg_media_volume_df.iloc[i,1]
            globals()[f"media{i+1}_volume"] = int(pg_media_volume_df.iloc[i,2])
            globals()[f"media{i+1}_share"] = f"{round(pg_media_volume_df.iloc[i,3])}%"
        except IndexError:
            globals()[f"media{i+1}_name"] = ""
            globals()[f"media{i+1}_volume"] = ""
            globals()[f"media{i+1}_share"] = ""

    # Update icons and metric labels
    channel = 0
    for icon in ["ch1_icon", "ch2_icon", "ch3_icon", "ch4_icon"]:
        shape = find_shape_by_name(prs.slides[4], icon)
        try:
            page_type = pg_media_volume_df.iloc[channel,0]
            slide_part, rId = shape.part, shape._element.blip_rId
            (new_impart_part, new_rid) = slide_part.get_or_add_image_part(f"img/{page_type}_icon.png")
            blip_fill = shape._element.blipFill
            blip = blip_fill.find(qn('a:blip'))
            blip.set(qn('r:embed'), new_rid)
            channel += 1
        except IndexError:
            sp = shape._sp
            sp.getparent().remove(sp)

    top_channel_metrics = {"ch1p": str(media1_share),
                    "ch1v": str(f"{media1_volume:,}"),
                    "ch1_label": "" if media1_name == 'X' else media1_name,
                    "ch2p": str(media2_share),
                    "ch2v": str(f"{media2_volume:,}"),
                    "ch2_label": "" if media2_name == 'X' else media2_name,
                    "ch3p": str(media3_share),
                    "ch3v": str(f"{media3_volume:,}"),
                    "ch3_label": "" if media3_name == 'X' else media3_name,
                    "ch4p": str(media4_share),
                    "ch4v": str(f"{media4_volume:,}"),
                    "ch4_label": "" if media4_name == 'X' else media4_name,
                    }

    for key, value in top_channel_metrics.items():
        replace_text(prs.slides[4], key, value)


    ## SLIDE 6(5) - BIG P&G DAILY SENTIMENT CHART

    pg_daily_sent_big = find_shape_by_name(prs.slides[5], "pg_daily_sent_big").chart
    pg_daily_sent_big.replace_data(pg_daily_sent_data)


    ## SLIDE 7(6) - COMPETITIVE SOV PIE CHART

    comp_sov_pie2_data = CategoryChartData()
    comp_sov_pie2_data.categories = brands_volume_df['shortName'].to_list()
    comp_sov_pie2_data.add_series("mentions", brands_volume_df['total'].to_list())

    comp_sov_pie_two = find_shape_by_name(prs.slides[6], "comp_sov_pie_two").chart
    comp_sov_pie_two.replace_data(comp_sov_pie2_data)


    ## SLIDE 7(6) - COMPETITIVE SENTIMENT BAR CHART

    comp_sentiment_bar_data = CategoryChartData()
    comp_sentiment_bar_data.categories = brands_volume_df['shortName'].to_list()
    comp_sentiment_bar_data.add_series("neutral", brands_volume_df['neu%'].to_list())
    comp_sentiment_bar_data.add_series("positive", brands_volume_df['pos%'].to_list())
    comp_sentiment_bar_data.add_series("negative", brands_volume_df['neg%'].to_list())

    comp_sentiment_bars = find_shape_by_name(prs.slides[6], "comp_sentiment_bars").chart
    comp_sentiment_bars.replace_data(comp_sentiment_bar_data)


    ## SLIDE 7(6) - TOP FOUR BRANDS VOLUME & PREV. PERIOD COMPARISONS

    top_brand_metrics = {"brand1_name": brand1_name,
                    "brand1_vol": f"{brand1_volume:,}",
                    "brand1_vs_prev": f"{abs(brand1_chng)}% vs. {abbreviate_number(brand1_prev)} previous period",
                    "brand2_name": brand2_name,
                    "brand2_vol": f"{brand2_volume:,}",
                    "brand2_vs_prev": f"{abs(brand2_chng)}% vs. {abbreviate_number(brand2_prev)} previous period",
                    "brand3_name": brand3_name,
                    "brand3_vol": f"{brand3_volume:,}",
                    "brand3_vs_prev": f"{abs(brand3_chng)}% vs. {abbreviate_number(brand3_prev)} previous period",
                    "brand4_name": brand4_name,
                    "brand4_vol": f"{brand4_volume:,}",
                    "brand4_vs_prev": f"{abs(brand4_chng)}% vs. {abbreviate_number(brand4_prev)} previous period",
                    }

    for key, value in top_brand_metrics.items():
        replace_text(prs.slides[6], key, value)


    brand1_arrow = find_shape_by_name(prs.slides[6], "brand1_arrow")
    if brand1_chng > 0:
        brand1_arrow.fill.fore_color.rgb = RGBColor(151, 215, 0)
        brand1_arrow.rotation = 180
    elif brand1_chng < 0:
        brand1_arrow.fill.fore_color.rgb = RGBColor(255, 109, 86)
        brand1_arrow.rotation = 0
    else:
        brand1_arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)

    brand2_arrow = find_shape_by_name(prs.slides[6], "brand2_arrow")
    if brand2_chng > 0:
        brand2_arrow.fill.fore_color.rgb = RGBColor(151, 215, 0)
        brand2_arrow.rotation = 180
    elif brand2_chng < 0:
        brand2_arrow.fill.fore_color.rgb = RGBColor(255, 109, 86)
        brand2_arrow.rotation = 0
    else:
        brand2_arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)

    brand3_arrow = find_shape_by_name(prs.slides[6], "brand3_arrow")
    if brand3_chng > 0:
        brand3_arrow.fill.fore_color.rgb = RGBColor(151, 215, 0)
        brand3_arrow.rotation = 180
    elif brand3_chng < 0:
        brand3_arrow.fill.fore_color.rgb = RGBColor(255, 109, 86)
        brand3_arrow.rotation = 0
    else:
        brand3_arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)

    brand4_arrow = find_shape_by_name(prs.slides[6], "brand4_arrow")
    if brand4_chng > 0:
        brand4_arrow.fill.fore_color.rgb = RGBColor(151, 215, 0)
        brand4_arrow.rotation = 180
    elif brand4_chng < 0:
        brand4_arrow.fill.fore_color.rgb = RGBColor(255, 109, 86)
        brand4_arrow.rotation = 0
    else:
        brand4_arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)


    ## DELETE COUNTRY SUMMARY SLIDE UNLESS REGION = EU OVERALL
    if region != "EU Overall":
        delete_slide(prs, 3)

    ## SAVE/RETURN THE FINAL REPORT PPTX
    # prs.save(f"Procter Gamble Corp. {region} - {fiscal_period}.pptx")

    binary_output = BytesIO()
    prs.save(binary_output)
    report_file = binary_output.getvalue()
    return report_file
