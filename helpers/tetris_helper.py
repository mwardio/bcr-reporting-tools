
from pptx import Presentation

'''
def print_layout_index_map(pptx_file):
    prs = Presentation(pptx_file)
    
    print("Slide Layout Index Map:\n")
    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name if hasattr(layout, "name") else "Unknown"
        print(f"Index {i}: {name}")

# Example usage:
print_layout_index_map("template_ynwa.pptx")

def debug_placeholders(slide):
    print("Placeholders on this slide:\n")
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        print(f"Index {phf.idx}: '{shape.name}' — text: '{shape.text.strip()}'")



prs = Presentation("template_ynwa.pptx")
# or a new one
slide = prs.slides.add_slide(prs.slide_layouts[12])  # Title slide layout

debug_placeholders(slide)
'''


### 1) Download a series of charts in BCR, joining them together

## Functions: fetch chart, format chart

## Inputs: List with ids (call for now), name of brand

### 2) Create a power point deck, with the chart and title and placeholders

# Slide 1 will need titles (so we will run a loop to add them)
# Slide 2 will need the raw bcr data, so that pptx can convert them to pie charts


# Cision logo?

### 3) Save to a file



### NOTE

## have two functions, one to create a chart, and one to create slide

## in the slide function, invoke the chart function for each chart, and the chart function should add the chart to the slide

## then save the slide

## user that logic, but your code that you have in one go right now


import requests as r
import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt


from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_LEGEND_POSITION


from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO



def gen_pptx_from_bcr_xlsx(bcr_token, start_date, end_date, excel_file_name):

    valid_sources = ['blog', 'bluesky', 'facebook', 'facebook_public', 'forum', 'instagram', 'instagram_public', 'linkedin', 'news', 'qq', 'reddit', 'review', 'threads', 'tiktok', 'tumblr', 'twitter', 'youtube', 'broadcast', 'podcast', 'print']

    def intro_add_picture(path, left, top, width):
        # Add picture
        img_path = path  # replace with your image file
        img_left = Inches(left)
        img_top = Inches(top)
        img_width = Inches(width)   # scale image width
        picture = slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
        return picture



    def fetch_data(bcr_token, project_id, metric, dim1, dim2, data_source_type, query_id, parent_categories, locations, content_sources, x_parent_categories, start_date, end_date, dim2Args:list):
        if dim2 == 'pageTypes':
            dim2Args = valid_sources
        auth = {'authorization': f'bearer {bcr_token}'}
        url = f"https://api.brandwatch.com/projects/{project_id}/data/{metric}/{dim1}/{dim2}"
        params = {'anyParentCategory': parent_categories,
                'startDate': start_date,
                'endDate': end_date,
                data_source_type: query_id,
                'location': locations,
                'pageType': content_sources,
                'xparentCategory': x_parent_categories,
                'dim2Args': dim2Args,
                'pageType': valid_sources
                }
        actual_call = r.get(url, headers = auth, params=params)
        full_url = actual_call.request.url
        debug_call = actual_call.json()
        if dim2 == 'pageTypes':
            results = [x for x in actual_call.json()['results']]
            results[0]['values'] = [x for x in results[0]['values'] if x['id'] in valid_sources]
            return results
        else:
            results = [x for x in actual_call.json()['results']]
            return results


    def gen_chart(companies_data, slide, top_offset):
        #print (len(companies_data))
        #section_width = Inches(9) / len(companies_data)
        
        for i, (company, data) in enumerate(companies_data.items()):
            section_width = Inches(data['section_width'])
            # X-position for this section
            x_offset = Inches(0.5) + i * section_width
            y_offset = top_offset
            
            # 1) Title (Company name)
            title_box = slide.shapes.add_textbox(x_offset, y_offset, section_width, Inches(0.4))
            title_tf = title_box.text_frame
            p = title_tf.paragraphs[0]
            run = p.add_run()
            run.text = company
            run.font.size = Pt(14)
            
            # 2) Pie chart (company-specific raw numbers)
            chart_data = CategoryChartData()
            chart_data.categories = data['date']
            chart_data.add_series('Volume', data['values'])

            if data['chart_type'] == 'volume over time':
                chart_type = XL_CHART_TYPE.LINE
            if data['chart_type'] == 'pie chart':
                chart_type = XL_CHART_TYPE.PIE
            
            chart = slide.shapes.add_chart(
                chart_type,
                #XL_CHART_TYPE.LINE,
                x_offset, y_offset + Inches(0.5),
                section_width - Inches(0.2), Inches(2.2),
                chart_data
            ).chart

            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(4)  # smaller font
            chart.chart_title.text_frame.paragraphs[0].font.bold = True 
            
            # Enable % data labels (normalized)
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.number_format = '0%'
            data_labels.font.size = Pt(10)
            
            # Add legend with all categories
            chart.has_legend = True
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            
            # 3) Text box (analysis placeholder)
            analysis_box = slide.shapes.add_textbox(
                x_offset, y_offset + Inches(2.8), section_width - Inches(0.2), Inches(2)
            )
            analysis_tf = analysis_box.text_frame
            analysis_tf.text = f"[analysis_for_{company}]"
        
        # Access X-axis (Category Axis)
            if data['chart_type'] == 'volume over time': 
                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(10)  # X-axis font size
                category_axis.tick_labels.font.bold = False

                # Access Y-axis (Value Axis)
                value_axis = chart.value_axis
                value_axis.tick_labels.font.size = Pt(10)  # Y-axis font size
                value_axis.tick_labels.font.bold = False
        return chart

        
    def gen_slide(prs, companies_data, top_offset, slide_name):
        slide_layout = prs.slide_layouts[12]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add main title
        #main_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
        #main_title_tf = main_title_box.text_frame
        #p = main_title_tf.paragraphs[0]
        #run = p.add_run()
        #run.text = "Main Overview"
        #run.font.size = Pt(28)
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_name
        # Define section widths

        

        #section_width = Inches(3)

        # Each company with multiple competitors (fictional real names + raw data)



        # Vertical offset for layout
        

        gen_chart(companies_data, slide, top_offset)

        return slide







    #bcr_token = '87670e72-2dfc-4280-84bd-3ada69d6b80b'
    #excel_file = 'jetblue.xlsx'


    logos = [{'name': 'main', 'path':'main_frame.png', 'left':3.5, 'top': 1, 'width': 6},
            {'name': 'left_logo_name', 'path':'left_logo_name.png', 'left':0.5, 'top': 0.5, 'width': 2}
            ]




    ### Intro Slide first


    prs = Presentation("./templates/template_ynwa.pptx")
    #remove_all_slides(prs)  

    # Add a title slide
    slide_layout = prs.slide_layouts[11]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = "Welcome to the JetBlue Media Report"
    title2_placeholder = slide.placeholders[1]
    title2_placeholder.text = "July 2025"
    #subtitle_placeholder = slide.placeholders[2]
    #subtitle_placeholder.text = "Powered by Brandwatch"


    # Add logos

    '''
    for logo in logos:
        intro_add_picture(logo['path'], logo['left'], logo['top'], logo['width'])
    '''
    '''

    # Get slide dimensions in order to calculate the center later
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Eneter desired textbox size
    text_width = Inches(3)
    text_height = Inches(2)

    # Left position (small gap from the left edge)
    text_left = Inches(0.5)

    # Top position: center vertically
    text_top = (slide_height - text_height) / 2

    # Add textbox
    textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = textbox.text_frame

    # Add text


    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Media Analysis Report - UK"
    run.font.size = Pt(20)
    run.font.bold = True


    '''




    # Step 1: Check all sheet names
    xls = pd.ExcelFile(excel_file_name)
    sheet_names = xls.sheet_names
    print("Available sheets:", sheet_names)




    for sheet in sheet_names:
        ### Create the graphs, and save them
        ### Use the slides in the graph
        file = pd.read_excel(excel_file_name, sheet_name=sheet)
        companies = file.to_dict(orient='records')
        for company in companies:
            company['data_source_ids'] = [company['data_source_ids']]
            if not pd.isna(company['dim2Args']):
                company['dim2Args'] = [company['dim2Args']]
                #company['dim2Args'] = list(map(int, company['dim2Args'].split(',')))
            elif pd.isna(company['dim2Args']):
                company['dim2Args'] = None

            #company['dim2Args'] = list(map(int, company['dim2Args'].split(',')))

            #company['data_source_ids'] = list(map(int, company['data_source_ids'].split(',')))

            #company['location'] = [name.strip() for name in company['location'].split(",")]
            #company['content_sources'] = [name.strip() for name in company['content_sources'].split(",")]
            #company['dim1Args'] = list(map(int, company['dim2Args'].split(','))) #TODO account for missing values


        chart_section_width = 10 /len(companies)

        companies_data = {}
        ## In this section, we will compile the data from bcr for each chart, resulting in the format required by pptx. However, we will also add keys for valuable information related to the chart, such as the type, fonts, etc
        for company in companies:
            company_chart_bcr = fetch_data(bcr_token, company['project_id'], company['metric'], company['dim1'], company['dim2'], company['data_source_type'], company['data_source_ids'], None, None, None, None, start_date, end_date, company['dim2Args'])
            
            company_data_for_ppt = {company['title']:{"date":[x["name"] for x in company_chart_bcr[0]['values']], 'values': [x['value'] for x in company_chart_bcr[0]['values']]}}
            company_data_for_ppt[company['title']]['chart_type'] = company['name']
            company_data_for_ppt[company['title']]['breakdown'] = company['dim2']
            company_data_for_ppt[company['title']]['section_width'] = chart_section_width
            companies_data.update(company_data_for_ppt)
            print ('test')
        #section_width = Inches(3)
        top_offset = Inches(1.2)
        slide_name = sheet
        gen_slide(prs, companies_data, top_offset, slide_name)

    #prs.save("companies_sov_slide_bcr_data_ex24.pptx")
    print("#YNWA")
    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)  # reset pointer
    return pptx_bytes
    


#2 is content simple

# 3 is title (with title, title 2 and subtitle)